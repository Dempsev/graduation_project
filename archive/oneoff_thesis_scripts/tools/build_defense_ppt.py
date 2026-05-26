from __future__ import annotations

import math
import os
import re
from pathlib import Path

import pandas as pd
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = Path(r"D:\graduation_project\答辩4.pptx")
OUT_DIR = ROOT / "output" / "defense_ppt"
ASSET_DIR = ROOT / "output" / "defense_ppt_assets"
PREVIEW_DIR = OUT_DIR / "preview_png"
OUT_PPTX = OUT_DIR / "苏炫丞_毕业答辩PPT_v1.pptx"
CONTACT_SHEET = OUT_DIR / "苏炫丞_毕业答辩PPT_v1_contact_sheet.png"


WIDE_W = 13.333333
WIDE_H = 7.5


COLORS = {
    "ink": RGBColor(35, 43, 52),
    "muted": RGBColor(95, 106, 117),
    "light": RGBColor(247, 248, 248),
    "line": RGBColor(214, 218, 220),
    "green": RGBColor(52, 145, 117),
    "blue": RGBColor(48, 86, 142),
    "orange": RGBColor(216, 129, 54),
    "gold": RGBColor(183, 139, 43),
    "red": RGBColor(176, 68, 55),
    "white": RGBColor(255, 255, 255),
}


def emu(n: float):
    return Inches(n)


def rgb_hex(c: RGBColor) -> str:
    return f"#{c[0]:02x}{c[1]:02x}{c[2]:02x}"


def delete_all_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001
    for sld_id in list(sld_id_lst):
        r_id = sld_id.rId
        prs.part.drop_rel(r_id)
        sld_id_lst.remove(sld_id)


def set_slide_bg(slide, color=COLORS["white"]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 24,
    bold: bool = False,
    color=COLORS["ink"],
    align=PP_ALIGN.LEFT,
    font="Microsoft YaHei",
    line_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_body_lines(slide, lines: list[str], x: float, y: float, w: float, h: float, size=21, color=COLORS["ink"]):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.line_spacing = 1.1
    return box


def add_title(slide, idx: int, title: str, kicker: str = "毕业答辩") -> None:
    add_textbox(slide, f"{idx:02d}", 0.56, 0.37, 0.56, 0.3, 12, True, COLORS["green"], PP_ALIGN.LEFT)
    add_textbox(slide, kicker, 1.08, 0.35, 1.2, 0.3, 10, False, COLORS["muted"], PP_ALIGN.LEFT)
    add_textbox(slide, title, 0.78, 0.72, 8.9, 0.55, 24, True, COLORS["ink"], PP_ALIGN.LEFT)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0.78), emu(1.34), emu(1.1), emu(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["green"]
    line.line.fill.background()
    footer = slide.shapes.add_textbox(emu(11.25), emu(7.05), emu(1.45), emu(0.22))
    footer.text_frame.text = "重庆大学航空航天学院"
    footer.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
    footer.text_frame.paragraphs[0].font.size = Pt(8)
    footer.text_frame.paragraphs[0].font.color.rgb = COLORS["muted"]


def add_chip(slide, text: str, x: float, y: float, w: float, color=COLORS["green"]):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, emu(x), emu(y), emu(w), emu(0.38))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.text = text
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS["white"]
    return shp


def add_round_rect(slide, x, y, w, h, fill=COLORS["light"], line=COLORS["line"], radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, emu(x), emu(y), emu(w), emu(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1)
    return shp


def add_card_text(slide, title, body, x, y, w, h, accent=COLORS["green"]):
    add_round_rect(slide, x, y, w, h, RGBColor(250, 251, 251), RGBColor(225, 229, 231))
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(x), emu(y), emu(0.08), emu(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_textbox(slide, title, x + 0.22, y + 0.18, w - 0.36, 0.35, 15, True, COLORS["ink"])
    add_textbox(slide, body, x + 0.22, y + 0.58, w - 0.36, h - 0.7, 12, False, COLORS["muted"])


def fit_image(slide, path: Path, x: float, y: float, w: float, h: float, border=True):
    if not path.exists():
        add_round_rect(slide, x, y, w, h)
        add_textbox(slide, f"缺少图片\n{path.name}", x + 0.2, y + 0.2, w - 0.4, h - 0.4, 14, False, COLORS["red"], PP_ALIGN.CENTER)
        return None
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = min(w / iw, h / ih)
    ww = iw * ratio
    hh = ih * ratio
    left = x + (w - ww) / 2
    top = y + (h - hh) / 2
    if border:
        add_round_rect(slide, x, y, w, h, RGBColor(255, 255, 255), RGBColor(226, 230, 232))
    return slide.shapes.add_picture(str(path), emu(left), emu(top), emu(ww), emu(hh))


def cover_image(slide, path: Path, x: float, y: float, w: float, h: float, opacity_overlay=True):
    if not path.exists():
        return
    slide.shapes.add_picture(str(path), emu(x), emu(y), emu(w), emu(h))
    if opacity_overlay:
        overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(x), emu(y), emu(w), emu(h))
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
        overlay.fill.transparency = 18
        overlay.line.fill.background()


def add_arrow(slide, x1, y1, x2, y2, color=COLORS["muted"], width=1.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, emu(x1), emu(y1), emu(x2), emu(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    try:
        c.line.end_arrowhead = True
    except Exception:
        pass
    return c


def add_process(slide, labels: list[str], x: float, y: float, w: float, h: float, colors=None):
    colors = colors or [COLORS["green"], COLORS["blue"], COLORS["orange"], COLORS["gold"]]
    gap = 0.2
    box_w = (w - gap * (len(labels) - 1)) / len(labels)
    for i, label in enumerate(labels):
        bx = x + i * (box_w + gap)
        add_round_rect(slide, bx, y, box_w, h, RGBColor(250, 251, 251), RGBColor(218, 222, 224))
        add_chip(slide, f"{i+1}", bx + 0.18, y + 0.18, 0.46, colors[i % len(colors)])
        add_textbox(slide, label, bx + 0.22, y + 0.72, box_w - 0.44, h - 0.86, 14, True, COLORS["ink"], PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            add_arrow(slide, bx + box_w + 0.03, y + h / 2, bx + box_w + gap - 0.03, y + h / 2)


def add_native_bar_chart(slide, categories, series, x, y, w, h, value_axis_max=None, legend=True):
    data = CategoryChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, emu(x), emu(y), emu(w), emu(h), data).chart
    chart.has_title = False
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.size = Pt(9)
    if value_axis_max:
        chart.value_axis.maximum_scale = value_axis_max
    for plot in chart.plots:
        plot.has_data_labels = True
        plot.data_labels.number_format = "0.0"
        plot.data_labels.font.size = Pt(8)
    palette = [COLORS["green"], COLORS["orange"], COLORS["blue"], COLORS["gold"]]
    for i, s in enumerate(chart.series):
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = palette[i % len(palette)]
        s.format.line.color.rgb = RGBColor(255, 255, 255)
    return chart


def add_target_bands(slide, x, y, w, h):
    bands = [(140, 180), (160, 200), (180, 220), (200, 240), (220, 260), (240, 280)]
    min_f, max_f = 130, 290
    axis_y = y + h - 0.32
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(x), emu(axis_y), emu(w), emu(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()
    for i, (lo, hi) in enumerate(bands):
        yy = y + i * 0.48
        left = x + (lo - min_f) / (max_f - min_f) * w
        width = (hi - lo) / (max_f - min_f) * w
        color = [COLORS["green"], COLORS["blue"], COLORS["orange"], COLORS["gold"], COLORS["red"], RGBColor(113, 92, 154)][i]
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, emu(left), emu(yy), emu(width), emu(0.28))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        add_textbox(slide, f"{lo}-{hi} Hz", x, yy - 0.02, 1.3, 0.24, 9, True, COLORS["ink"])
    for f in [140, 180, 220, 260, 280]:
        xx = x + (f - min_f) / (max_f - min_f) * w
        tick = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(xx), emu(axis_y - 0.06), emu(0.012), emu(0.13))
        tick.fill.solid()
        tick.fill.fore_color.rgb = COLORS["muted"]
        tick.line.fill.background()
        add_textbox(slide, str(f), xx - 0.15, axis_y + 0.08, 0.35, 0.18, 7, False, COLORS["muted"], PP_ALIGN.CENTER)


def add_formula_overlap(slide, x, y):
    add_round_rect(slide, x, y, 5.5, 1.15, RGBColor(253, 250, 241), RGBColor(226, 208, 156))
    add_textbox(slide, "O = max(0, min(fu, gu) - max(fl, gl))", x + 0.25, y + 0.16, 5.0, 0.34, 19, True, COLORS["ink"], PP_ALIGN.CENTER, font="Cambria Math")
    add_textbox(slide, "O 表示真实带隙与目标频带的重叠宽度，覆盖率 = O / 目标频带宽度。", x + 0.35, y + 0.62, 4.8, 0.28, 11, False, COLORS["muted"], PP_ALIGN.CENTER)


def load_tables():
    sixband = pd.read_csv(ROOT / "data/analysis/thesis_ch5_titleless_cn_bundle_v1/sixband_predictor_vs_ga_summary_v1.csv")
    three = pd.read_csv(ROOT / "data/analysis/targetband_four_arm_baseline_v10_fullpool_v1/targetband_three_method_plot_values_v10_cn.csv")
    return sixband, three


def build_preview_images(slide_texts: list[tuple[str, str]], pptx_path: Path) -> None:
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    width, height = 1280, 720
    try:
        font_title = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 34)
        font_body = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 19)
        font_small = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 14)
    except Exception:
        font_title = font_body = font_small = ImageFont.load_default()
    preview_paths = []
    for i, (title, body) in enumerate(slide_texts, 1):
        im = Image.new("RGB", (width, height), "#f8f9f8")
        draw = ImageDraw.Draw(im)
        draw.rectangle([0, 0, width, 54], fill="#eef4f1")
        draw.rectangle([58, 86, 170, 92], fill="#349175")
        draw.text((60, 28), f"{i:02d}", fill="#349175", font=font_small)
        draw.text((60, 100), title, fill="#232b34", font=font_title)
        y = 172
        for line in body.split("\n")[:8]:
            draw.text((82, y), line, fill="#53616e", font=font_body)
            y += 34
        draw.text((width - 350, height - 35), f"preview: {pptx_path.name}", fill="#9aa3aa", font=font_small)
        out = PREVIEW_DIR / f"slide_{i:02d}.png"
        im.save(out)
        preview_paths.append(out)
    thumbs = []
    for p in preview_paths:
        im = Image.open(p).resize((320, 180))
        thumbs.append(im)
    sheet = Image.new("RGB", (320 * 3, 180 * math.ceil(len(thumbs) / 3)), "white")
    for i, im in enumerate(thumbs):
        sheet.paste(im, ((i % 3) * 320, (i // 3) * 180))
    sheet.save(CONTACT_SHEET)


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(TEMPLATE))
    prs.slide_width = emu(WIDE_W)
    prs.slide_height = emu(WIDE_H)
    delete_all_slides(prs)
    blank = prs.slide_layouts[1]
    sixband, three = load_tables()

    assets = {
        "cover": ASSET_DIR / "template_media" / "image5.png",
        "unit": ROOT / "data/analysis/thesis_ch2_v1/figures/ch2_snake_fourier_overlay_mesh_ep100_step18.png",
        "workflow_geom": ROOT / "data/analysis/thesis_ch2_v1/figures/figure_2_x_snake_fourier_overlay_workflow_v1.png",
        "problem": ROOT / "data/analysis/thesis_ch2_v1/figures/figure_2_1_problem_boundary.png",
        "ga_conv": ROOT / "data/analysis/thesis_ch5_titleless_cn_bundle_v1/figures/figure_5_7a_multiband_ga_convergence_cn_titleless.png",
        "ga_best": ROOT / "data/analysis/thesis_ch5_titleless_cn_bundle_v1/figures/figure_5_7b_multiband_ga_final_best_cn_titleless.png",
        "candidate": ROOT / "data/analysis/thesis_ch5_titleless_cn_bundle_v1/figures/figure_5_7d_candidate_pool_predicted_best_cn_titleless.png",
        "three_method": ROOT / "data/analysis/targetband_four_arm_baseline_v10_fullpool_v1/figures/figure_5_6a_target_overlap_comparison_v10_cn.png",
        "disp_180": ROOT / "data/analysis/canonical_local_robustness_v1/dispersion_plots/band180_220_ep248_dispersion_compare_v1.png",
        "disp_200": ROOT / "data/analysis/canonical_local_robustness_v1/dispersion_plots/band200_240_ep193_dispersion_compare_v1.png",
        "disp_240": ROOT / "data/analysis/canonical_local_robustness_v1/dispersion_plots/band240_280_ep253_dispersion_compare_v1.png",
        "mode_240": ROOT / "data/analysis/canonical_mode_shapes_v1/band240_280_ep253/stage4_validation_targetband_local_robustness_v1_band240_280_ep253_center_lower_edge.png",
    }

    preview_texts: list[tuple[str, str]] = []

    # 1
    slide = prs.slides.add_slide(blank); set_slide_bg(slide)
    cover_image(slide, assets["cover"], 0, 0, WIDE_W, WIDE_H, True)
    add_textbox(slide, "基于物理-数据协同优化算法的\n一种力学超结构设计", 0.88, 1.25, 7.4, 1.5, 33, True, COLORS["ink"])
    add_textbox(slide, "毕业答辩", 0.92, 0.78, 1.5, 0.28, 14, True, COLORS["green"])
    add_textbox(slide, "姓名：苏炫丞    专业：工程力学\n指导教师：XXX\n重庆大学航空航天学院", 0.94, 4.5, 4.4, 0.9, 15, False, COLORS["muted"])
    preview_texts.append(("封面", "基于物理-数据协同优化算法的一种力学超结构设计\n姓名、专业、导师、学院"))

    # 2
    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 2, "周期结构带隙调控及工程应用背景")
    fit_image(slide, assets["problem"], 7.25, 1.55, 4.95, 4.6)
    add_textbox(slide, "研究对象不是“随便找带隙”，而是让带隙服务于指定频率范围。", 0.9, 1.62, 5.75, 0.78, 25, True, COLORS["ink"])
    add_body_lines(slide, ["周期结构可通过单胞重复调控弹性波传播", "在特定频段形成带隙，可用于振动隔离与噪声控制", "工程设计更关心目标频带内是否真正形成有效带隙"], 1.05, 2.72, 5.35, 1.72, 18)
    add_target_bands(slide, 1.05, 5.0, 5.35, 1.6)
    preview_texts.append(("研究背景", "声子晶体/周期结构调控弹性波\n工程更关心目标频带内是否有有效带隙"))

    # 3
    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 3, "目标频带设计中的计算成本问题")
    add_textbox(slide, "传统搜索的瓶颈：每个候选都要进入高成本有限元验证。", 0.92, 1.55, 6.4, 0.55, 24, True)
    add_process(slide, ["几何建模", "网格划分", "COMSOL\n特征频率求解", "频散后处理", "带隙与目标频带\n重叠评价"], 0.9, 2.45, 11.45, 1.45, [COLORS["green"], COLORS["blue"], COLORS["orange"], COLORS["gold"], COLORS["red"]])
    add_card_text(slide, "候选空间大", "结构族与连续参数共同决定候选，组合数量快速增加。", 1.02, 4.65, 2.6, 1.25, COLORS["green"])
    add_card_text(slide, "仿真成本高", "频散计算需要建模、网格与多波矢特征值求解。", 3.95, 4.65, 2.6, 1.25, COLORS["blue"])
    add_card_text(slide, "无效候选多", "很多结构形成带隙，但位置不在目标频带内。", 6.88, 4.65, 2.6, 1.25, COLORS["orange"])
    add_card_text(slide, "预算被浪费", "盲目搜索会把 COMSOL 资源消耗在低价值候选上。", 9.81, 4.65, 2.6, 1.25, COLORS["red"])
    preview_texts.append(("研究问题", "每个候选都跑 COMSOL 很贵\n候选空间大、无效候选多、验证预算有限"))

    # 4
    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 4, "物理计算-机器学习预测-真实优化相结合的研究路线")
    add_textbox(slide, "COMSOL 提供物理真值，机器学习负责候选筛选，真实 GA 提供优化基准，最后用真实验证说明筛选是否有效。", 0.95, 1.48, 11.4, 0.68, 22, True)
    add_process(slide, ["物理真值计算\nCOMSOL 频散/带隙", "条件预测模型\n分类 + 回归", "真实 GA 优化\nCOMSOL overlap 适应度", "预测筛选对比\nTop5 / 随机 / GA"], 0.9, 2.82, 11.55, 1.85)
    add_round_rect(slide, 1.3, 5.35, 10.7, 0.78, RGBColor(242, 247, 245), RGBColor(205, 224, 216))
    add_textbox(slide, "主线概括：物理真值计算 → 条件预测建模 → 真实闭环优化 → 预测筛选对比", 1.55, 5.52, 10.2, 0.28, 19, True, COLORS["green"], PP_ALIGN.CENTER)
    preview_texts.append(("研究路线", "物理真值计算 → 条件预测建模 → 真实闭环优化 → 预测筛选对比"))

    # 5
    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 5, "二维参数化声子晶体及目标频带任务")
    fit_image(slide, assets["unit"], 0.9, 1.55, 4.9, 4.55)
    add_textbox(slide, "任务定义", 6.35, 1.6, 2.0, 0.4, 22, True)
    add_body_lines(slide, ["输入：结构参数、结构族信息、目标频带上下限", "输出：目标频带是否有效开启、目标频带覆盖率", "同一结构面对不同频带时，标签可以不同"], 6.35, 2.18, 5.7, 1.35, 17)
    add_textbox(slide, "六个目标频带", 6.35, 4.1, 2.0, 0.3, 16, True, COLORS["green"])
    add_target_bands(slide, 6.35, 4.55, 5.25, 1.7)
    preview_texts.append(("研究对象", "二维周期单胞；输入结构参数/结构族/目标频带；输出开启状态与覆盖率"))

    # 6
    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 6, "基于 Bloch 周期边界的频散计算方法")
    add_process(slide, ["单胞模型", "Bloch-Floquet\n周期边界", "波矢路径\n特征频率求解", "频散曲线", "带隙提取\n目标重叠"], 0.82, 1.75, 11.7, 1.35)
    add_formula_overlap(slide, 0.98, 3.75)
    fit_image(slide, assets["disp_180"], 6.85, 3.35, 5.4, 2.75)
    add_textbox(slide, "这一页只回答“怎么算”：真实带隙与目标频带的重叠宽度作为统一评价指标。", 1.03, 5.25, 5.35, 0.55, 16, False, COLORS["muted"])
    preview_texts.append(("COMSOL 频散计算", "Bloch 周期边界；沿波矢路径求特征频率；提取带隙并计算目标 overlap"))

    # 7
    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 7, "单胞几何模型与参数化描述")
    fit_image(slide, assets["workflow_geom"], 0.82, 1.55, 6.25, 4.9)
    add_card_text(slide, "参数化控制", "候选结构由结构族与连续几何参数共同确定。", 7.55, 1.65, 4.25, 1.05, COLORS["green"])
    add_card_text(slide, "有效性检查", "几何有效性、接触有效性和求解稳定性共同决定候选能否进入 COMSOL。", 7.55, 3.05, 4.25, 1.05, COLORS["blue"])
    add_card_text(slide, "优化惩罚", "无效结构在 GA 中给予惩罚，防止搜索被失败样本牵引。", 7.55, 4.45, 4.25, 1.05, COLORS["orange"])
    preview_texts.append(("几何参数化", "结构族 + 连续参数；几何/接触/求解稳定性检查；无效结构惩罚"))

    # 8
    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 8, "面向目标频带的条件分类与回归模型")
    add_textbox(slide, "把目标频带上下限也作为输入，模型回答的是“这个结构是否适合当前频段”。", 0.92, 1.55, 10.8, 0.55, 23, True)
    add_round_rect(slide, 1.0, 2.55, 3.15, 2.55, RGBColor(245, 250, 248), RGBColor(205, 224, 216))
    add_textbox(slide, "输入条件", 1.28, 2.84, 2.5, 0.35, 18, True, COLORS["green"], PP_ALIGN.CENTER)
    add_body_lines(slide, ["结构参数", "结构族信息", "目标频带 fl, fu"], 1.38, 3.42, 2.3, 1.0, 15, COLORS["ink"])
    add_arrow(slide, 4.22, 3.82, 5.05, 3.82, COLORS["green"], 2.0)
    add_round_rect(slide, 5.15, 2.55, 3.15, 2.55, RGBColor(248, 250, 253), RGBColor(204, 216, 235))
    add_textbox(slide, "条件预测模型", 5.43, 2.84, 2.5, 0.35, 18, True, COLORS["blue"], PP_ALIGN.CENTER)
    add_body_lines(slide, ["分类：是否开启", "回归：overlap / coverage", "用于快速排序"], 5.52, 3.42, 2.35, 1.0, 15, COLORS["ink"])
    add_arrow(slide, 8.35, 3.82, 9.15, 3.82, COLORS["blue"], 2.0)
    add_round_rect(slide, 9.25, 2.55, 3.15, 2.55, RGBColor(253, 248, 242), RGBColor(232, 211, 191))
    add_textbox(slide, "输出", 9.58, 2.84, 2.4, 0.35, 18, True, COLORS["orange"], PP_ALIGN.CENTER)
    add_body_lines(slide, ["有效候选概率", "重叠宽度预测", "验证优先级"], 9.65, 3.42, 2.25, 1.0, 15, COLORS["ink"])
    preview_texts.append(("机器学习条件预测", "同一结构面对不同目标频带标签不同；分类判断开启，回归预测覆盖程度"))

    # 9
    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 9, "预测模型的作用：候选筛选而非最终判据")
    add_textbox(slide, "预测模型不替代 COMSOL；它的价值在于把有限验证预算优先分配给更值得验证的结构。", 0.9, 1.55, 11.5, 0.7, 24, True)
    add_card_text(slide, "机器学习预测", "速度快；适合大规模初筛与排序；可信度需要真实频散验证。", 1.0, 2.75, 5.25, 1.7, COLORS["green"])
    add_card_text(slide, "COMSOL 频散计算", "计算慢；物理可信度高；作为标签来源、优化适应度和最终判据。", 7.05, 2.75, 5.25, 1.7, COLORS["blue"])
    add_round_rect(slide, 1.45, 5.25, 10.4, 0.72, RGBColor(252, 248, 239), RGBColor(230, 207, 149))
    add_textbox(slide, "答辩回答口径：不是“用模型代替物理”，而是“在相同 COMSOL 预算下更快找到有效结构”。", 1.75, 5.42, 9.8, 0.28, 17, True, COLORS["gold"], PP_ALIGN.CENTER)
    preview_texts.append(("模型定位", "预测做初筛与排序；最终性能以 COMSOL 真实频散为准；关注相同验证预算下的发现效率"))

    # 10
    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 10, "基于真实频散计算的目标频带遗传优化")
    add_process(slide, ["初始种群", "几何生成", "COMSOL 频散", "overlap 评价", "选择/交叉/变异", "下一代"], 0.82, 1.8, 11.8, 1.45, [COLORS["green"], COLORS["blue"], COLORS["orange"], COLORS["gold"], COLORS["red"], RGBColor(113, 92, 154)])
    add_textbox(slide, "关键点：GA 的适应度来自真实 COMSOL overlap，不是预测模型输出。", 1.02, 3.85, 7.0, 0.45, 23, True, COLORS["ink"])
    add_body_lines(slide, ["每个个体对应一组结构参数", "生成结构并进行有效性检查", "频散后处理提取目标频带重叠宽度", "六个目标频带分别优化，形成真实基准"], 1.05, 4.6, 5.4, 1.55, 17)
    fit_image(slide, assets["ga_conv"], 7.0, 3.55, 5.25, 2.8)
    preview_texts.append(("真实 COMSOL-GA", "每个个体跑 COMSOL；overlap 是适应度；真实 GA 提供优化基准"))

    # 11
    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 11, "不同目标频带下的真实优化结果")
    cats = list(sixband["band_label"])
    ga_vals = [float(v) for v in sixband["real_ga_12gen_best_comsol_overlap_Hz"]]
    add_native_bar_chart(slide, cats, [("真实 GA 最优 overlap / Hz", ga_vals)], 0.95, 1.72, 7.25, 4.55, 45, False)
    add_textbox(slide, "结果读法", 8.7, 1.82, 2.0, 0.35, 20, True, COLORS["green"])
    add_body_lines(slide, ["180-220 Hz 实现完整或近完整覆盖", "160-200 Hz 与 200-240 Hz 取得较高覆盖", "220-260 Hz 与 240-280 Hz 高频目标覆盖较低", "高频困难指向结构族和参数空间表达能力边界"], 8.7, 2.42, 3.3, 2.1, 16)
    fit_image(slide, assets["ga_best"], 8.55, 4.92, 3.55, 1.55, border=False)
    preview_texts.append(("GA 优化结果", "中频段覆盖较好；高频段弱；问题不只是模型，而是结构族/参数空间边界"))

    # 12
    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 12, "独立候选集验证：预测 Top5 vs 随机候选")
    add_textbox(slide, "验证逻辑：从独立候选集出发，比较相同 COMSOL 验证预算下谁更容易找到有效结构。", 0.92, 1.55, 11.0, 0.55, 22, True)
    add_process(slide, ["独立候选池", "排除训练集", "排除 GA20 历史", "排除已有候选", "COMSOL 真实验证"], 0.9, 2.5, 11.55, 1.35)
    add_card_text(slide, "比较对象", "预测 Top5 候选、随机候选、真实 GA 优化基准。", 1.05, 4.55, 3.3, 1.3, COLORS["green"])
    add_card_text(slide, "评价指标", "真实 target overlap、覆盖率、有效候选比例、预算内最优结果。", 5.0, 4.55, 3.3, 1.3, COLORS["blue"])
    add_card_text(slide, "防止复盘", "验证集不直接来自训练集或历史搜索记录，避免只证明模型记住旧样本。", 8.95, 4.55, 3.3, 1.3, COLORS["orange"])
    preview_texts.append(("预测筛选设置", "独立候选集；排除训练/GA20/已有候选；比较 Top5、随机和 GA"))

    # 13
    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 13, "预测筛选在多数频带下优于随机候选")
    top_vals = [float(v) for v in sixband["predictor_top1_comsol_truth_overlap_Hz"]]
    ga_vals = [float(v) for v in sixband["real_ga_12gen_best_comsol_overlap_Hz"]]
    add_native_bar_chart(slide, cats, [("预测 Top 候选 COMSOL 真值", top_vals), ("真实 GA 基准", ga_vals)], 0.8, 1.62, 7.55, 4.65, 45, True)
    add_textbox(slide, "核心发现", 8.82, 1.78, 2.0, 0.35, 20, True, COLORS["green"])
    add_body_lines(slide, ["六个目标频带中，预测 Top5 在五个频带上的最优 overlap 高于随机候选", "180-220 Hz 与 200-240 Hz 的预测结果接近 GA 基准", "说明模型能把高性能结构邻域排到验证队列前面"], 8.82, 2.35, 3.3, 1.6, 16)
    add_round_rect(slide, 8.82, 4.75, 3.4, 1.0, RGBColor(242, 247, 245), RGBColor(205, 224, 216))
    add_textbox(slide, "价值不是替代 GA，而是提高有限 COMSOL 预算下的候选发现效率。", 9.05, 5.02, 2.92, 0.36, 14, True, COLORS["green"], PP_ALIGN.CENTER)
    preview_texts.append(("Top5 对比", "预测 Top5 多数频带优于随机；180-220 和 200-240 接近 GA 基准"))

    # 14
    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 14, "典型目标频带结构频散验证")
    fit_image(slide, assets["disp_180"], 0.75, 1.55, 3.9, 3.3)
    fit_image(slide, assets["disp_200"], 4.72, 1.55, 3.9, 3.3)
    fit_image(slide, assets["disp_240"], 8.69, 1.55, 3.9, 3.3)
    add_card_text(slide, "180-220 Hz", "正例：预测筛选接近真实优化结果。", 0.85, 5.3, 3.55, 0.95, COLORS["green"])
    add_card_text(slide, "200-240 Hz", "过渡频带：仍保持较好的验证效果。", 4.82, 5.3, 3.55, 0.95, COLORS["blue"])
    add_card_text(slide, "240-280 Hz", "困难频带：预测、随机和 GA 均表现有限。", 8.79, 5.3, 3.55, 0.95, COLORS["red"])
    preview_texts.append(("典型案例", "180-220 正例；200-240 过渡；240-280 方法边界"))

    # 15
    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 15, "高频目标频带的困难与方法边界")
    add_textbox(slide, "高频频带表现弱不等于方法失效，而是揭示了当前结构族的适用边界。", 0.92, 1.55, 11.3, 0.62, 24, True)
    add_body_lines(slide, ["220-260 Hz 和 240-280 Hz 中，预测候选、随机候选和 GA 结果都有限", "问题不只是预测误差，也可能是结构族与参数空间难以表达高频带隙", "高频补样能找到部分弱频带候选，但局部模态特征不等于完整目标带隙", "最终仍需 COMSOL 频散验证"], 0.95, 2.48, 5.9, 2.1, 17)
    fit_image(slide, assets["mode_240"], 7.2, 2.18, 4.8, 3.55)
    add_round_rect(slide, 1.12, 5.62, 5.3, 0.65, RGBColor(252, 244, 242), RGBColor(228, 184, 177))
    add_textbox(slide, "答辩口径：承认边界，但把边界解释为结构表达能力问题，而不是预测框架整体失效。", 1.35, 5.78, 4.82, 0.28, 13, True, COLORS["red"], PP_ALIGN.CENTER)
    preview_texts.append(("高频困难", "高频目标覆盖有限；边界来自结构族/参数空间表达能力；仍以 COMSOL 验证为准"))

    # 16
    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 16, "主要结论")
    conclusions = [
        ("COMSOL 真值体系", "建立了目标频带设计的频散计算流程和 overlap/coverage 评价指标。", COLORS["green"]),
        ("条件预测模型", "构建了面向目标频带的分类与回归模型，用于候选筛选和排序。", COLORS["blue"]),
        ("真实优化基准", "建立了基于真实 COMSOL 频散计算的闭环 GA 优化基准。", COLORS["orange"]),
        ("筛选有效性", "独立验证表明预测 Top5 多数频带优于随机候选，提高验证预算效率。", COLORS["gold"]),
    ]
    for i, (t, b, c) in enumerate(conclusions):
        x = 0.95 + (i % 2) * 5.85
        y = 1.75 + (i // 2) * 2.05
        add_card_text(slide, t, b, x, y, 5.1, 1.45, c)
    add_round_rect(slide, 1.25, 6.0, 10.75, 0.58, RGBColor(242, 247, 245), RGBColor(205, 224, 216))
    add_textbox(slide, "一句话结论：机器学习预测不能替代 COMSOL，但可以提高 COMSOL 验证资源的使用效率。", 1.55, 6.15, 10.15, 0.22, 16, True, COLORS["green"], PP_ALIGN.CENTER)
    preview_texts.append(("主要结论", "COMSOL 真值体系；条件预测；真实 GA 基准；预测筛选提升预算效率"))

    # 17
    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 17, "不足与展望")
    add_card_text(slide, "结构表达能力", "当前结构族和参数化空间对高频目标频带表达能力有限。", 1.0, 1.9, 3.45, 2.0, COLORS["red"])
    add_card_text(slide, "数据分布", "训练数据在高频目标区间仍然稀疏，模型外推能力受限。", 4.95, 1.9, 3.45, 2.0, COLORS["orange"])
    add_card_text(slide, "后续方向", "引入更丰富结构族、主动采样策略和更强物理约束的预测模型。", 8.9, 1.9, 3.45, 2.0, COLORS["green"])
    add_textbox(slide, "后续研究将进一步提高高频目标频带样本覆盖度，并探索更具表达能力的结构参数化方法。", 1.28, 5.05, 10.85, 0.55, 23, True, COLORS["ink"], PP_ALIGN.CENTER)
    preview_texts.append(("不足与展望", "高频结构表达能力有限；高频数据稀疏；后续丰富结构族与主动采样"))

    # 18
    slide = prs.slides.add_slide(blank); set_slide_bg(slide)
    cover_image(slide, assets["cover"], 0, 0, WIDE_W, WIDE_H, True)
    add_textbox(slide, "请各位老师批评指正", 2.1, 2.7, 8.9, 0.75, 42, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_textbox(slide, "谢谢！", 5.55, 3.65, 2.2, 0.5, 25, True, COLORS["green"], PP_ALIGN.CENTER)
    preview_texts.append(("致谢", "请各位老师批评指正"))

    prs.save(str(OUT_PPTX))
    build_preview_images(preview_texts, OUT_PPTX)
    print(OUT_PPTX)
    print(CONTACT_SHEET)


if __name__ == "__main__":
    main()
