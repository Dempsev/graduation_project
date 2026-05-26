from __future__ import annotations

from pathlib import Path

import pandas as pd
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor

from build_defense_ppt import (
    ASSET_DIR,
    COLORS,
    TEMPLATE,
    WIDE_H,
    WIDE_W,
    add_body_lines,
    add_card_text,
    add_formula_overlap,
    add_native_bar_chart,
    add_process,
    add_round_rect,
    add_target_bands,
    add_textbox,
    add_title,
    cover_image,
    delete_all_slides,
    emu,
    fit_image,
    set_slide_bg,
)


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "output" / "defense_ppt"
PREVIEW_DIR = OUT_DIR / "preview_png_v2"
OUT_PPTX = OUT_DIR / "苏炫丞_毕业答辩PPT_v2_新大纲新图版.pptx"
CONTACT_SHEET = OUT_DIR / "苏炫丞_毕业答辩PPT_v2_新大纲新图版_contact_sheet.png"


def build_preview(slide_texts: list[tuple[str, str]]) -> None:
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    width, height = 1280, 720
    try:
        title_font = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 34)
        body_font = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 19)
        small_font = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 14)
    except Exception:
        title_font = body_font = small_font = ImageFont.load_default()
    paths = []
    for idx, (title, body) in enumerate(slide_texts, 1):
        im = Image.new("RGB", (width, height), "#f8f9f8")
        draw = ImageDraw.Draw(im)
        draw.rectangle([0, 0, width, 54], fill="#eef4f1")
        draw.rectangle([58, 86, 170, 92], fill="#349175")
        draw.text((60, 28), f"{idx:02d}", fill="#349175", font=small_font)
        draw.text((60, 100), title, fill="#232b34", font=title_font)
        y = 172
        for line in body.split("\n")[:8]:
            draw.text((82, y), line, fill="#53616e", font=body_font)
            y += 34
        out = PREVIEW_DIR / f"slide_{idx:02d}.png"
        im.save(out)
        paths.append(out)
    thumbs = [Image.open(p).resize((320, 180)) for p in paths]
    sheet = Image.new("RGB", (320 * 3, 180 * 6), "white")
    for i, thumb in enumerate(thumbs):
        sheet.paste(thumb, ((i % 3) * 320, (i // 3) * 180))
    sheet.save(CONTACT_SHEET)


def image_path(rel: str) -> Path:
    return ROOT / rel


def load_new_tables():
    ch4 = pd.read_csv(ROOT / "research_validation/ch4_ga_real_optimization/ch4_ga_summary_20gen.csv")
    strict = pd.read_csv(ROOT / "research_validation/ch5_strict_holdout_validation/ch5_strict_holdout_summary.csv")
    vs_ga = pd.read_csv(ROOT / "research_validation/ch5_strict_holdout_validation/ch5_strict_holdout_vs_ga20.csv")
    return ch4, strict, vs_ga


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(TEMPLATE))
    prs.slide_width = emu(WIDE_W)
    prs.slide_height = emu(WIDE_H)
    delete_all_slides(prs)
    blank = prs.slide_layouts[1]
    ch4, strict, vs_ga = load_new_tables()

    assets = {
        "cover": ASSET_DIR / "template_media" / "image5.png",
        "ch1_route": image_path("data/analysis/thesis_ch1_v1/figures/图片2.png"),
        "unit": image_path("data/analysis/thesis_ch2_v1/figures/ch2_snake_fourier_overlay_mesh_ep100_step18.png"),
        "geom_workflow": image_path("data/analysis/thesis_ch2_v1/figures/figure_2_x_snake_fourier_overlay_workflow_v1.png"),
        "overlap": image_path("data/analysis/thesis_ch2_v1/figures/图片3.png"),
        "mesh": image_path("data/research_validation/ch2_mesh_reliability_v1/figures/mesh_dispersion_overlay_v1.png"),
        "typical_disp": image_path("data/research_validation/ch2_typical_dispersion/figures/ch2_typical_local_dispersion_compare.png"),
        "ml_model": image_path("research_validation/ch3_predictor_v12_figures/ch3_model_structure.png"),
        "ml_metrics": image_path("research_validation/ch3_predictor_v12_figures/matlab_polished/ch3_matlab_topk_mean_cover_lines.png"),
        "ga_flow": image_path("research_validation/ch4_ga_real_optimization/figures/ch4_fig4_1_real_ga_flowchart.png"),
        "ga_comsol": image_path("research_validation/ch4_ga_real_optimization/figures/ch4_fig4_2_comsol_evaluation_flowchart.png"),
        "ga_conv": image_path("research_validation/ch4_ga_real_optimization/figures/ch4_fig4_3_ga_convergence_20gen_papercolor.png"),
        "ga_best": image_path("research_validation/ch4_ga_real_optimization/figures/ch4_fig4_4_best_overlap_bar_20gen_papercolor.png"),
        "ga_cells": image_path("research_validation/ch4_ga_real_optimization/figures/ch4_fig4_6_best_unit_cells_6bands_comsol.png"),
        "holdout_pipe": image_path("research_validation/ch5_strict_holdout_validation/figures/ch5_strict_fig1_holdout_pipeline.png"),
        "strict_overlap": image_path("research_validation/ch5_strict_holdout_validation/figures/ch5_strict_fig3_pred_vs_random_best_overlap.png"),
        "strict_cover": image_path("research_validation/ch5_strict_holdout_validation/figures/ch5_strict_fig4_pred_vs_random_best_cover.png"),
        "strict_ratio": image_path("research_validation/ch5_strict_holdout_validation/figures/ch5_strict_fig5_vs_ga20_ratio_notitle.png"),
        "typical_cells": image_path("research_validation/ch5_strict_holdout_validation/figures/ch5_strict_fig7_unit_cells_pred_vs_ga20_redraw.png"),
        "typical_disp_strict": image_path("research_validation/ch5_strict_holdout_validation/figures/ch5_strict_fig8_dispersion_pred_vs_ga20_redraw.png"),
        "highfreq": image_path("research_validation/ch5_prediction_vs_ga/figures/ch5_fig5_11_highfreq_boundary_analysis.png"),
    }

    preview: list[tuple[str, str]] = []

    slide = prs.slides.add_slide(blank); set_slide_bg(slide)
    cover_image(slide, assets["cover"], 0, 0, WIDE_W, WIDE_H, True)
    add_textbox(slide, "基于物理-数据协同优化算法的\n一种力学超结构设计", 0.88, 1.2, 7.6, 1.48, 33, True, COLORS["ink"])
    add_textbox(slide, "毕业答辩", 0.92, 0.78, 1.5, 0.28, 14, True, COLORS["green"])
    add_textbox(slide, "姓名：苏炫丞    专业：工程力学\n指导教师：XXX\n重庆大学航空航天学院", 0.94, 4.5, 4.4, 0.9, 15, False, COLORS["muted"])
    preview.append(("封面", "题目、姓名、专业、导师、学院；保持干净。"))

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 2, "周期结构带隙调控及工程应用背景")
    fit_image(slide, assets["typical_disp"], 7.0, 1.48, 5.35, 4.8)
    add_textbox(slide, "工程设计关注的是：指定频率范围内是否出现可用带隙。", 0.9, 1.58, 5.6, 0.62, 24, True)
    add_body_lines(slide, ["声子晶体/周期结构可调控弹性波传播", "特定频段形成带隙，可服务振动隔离、噪声控制", "目标频带设计同时要求带隙宽度与位置匹配"], 1.03, 2.55, 5.3, 1.5, 18)
    add_target_bands(slide, 1.0, 4.85, 5.35, 1.6)
    preview.append(("研究背景", "声子晶体可调控弹性波；工程更关心指定频带内的带隙。"))

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 3, "目标频带设计中的计算成本问题")
    add_textbox(slide, "直接参数扫描或盲目优化，会把大量 COMSOL 预算花在无效候选上。", 0.92, 1.55, 11.0, 0.55, 23, True)
    add_process(slide, ["几何建模", "网格划分", "COMSOL\n特征频率求解", "频散曲线后处理", "带隙提取\n目标评价"], 0.88, 2.42, 11.55, 1.45)
    fit_image(slide, assets["mesh"], 0.95, 4.45, 5.3, 1.85)
    add_card_text(slide, "核心矛盾", "候选结构空间大，真实频散计算成本高，而许多候选对目标频带没有贡献。", 6.85, 4.55, 4.95, 1.45, COLORS["orange"])
    preview.append(("研究问题", "每个候选都跑 COMSOL：候选空间大、仿真成本高、无效样本多。"))

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 4, "物理计算-机器学习预测-真实优化相结合的研究路线")
    add_textbox(slide, "COMSOL 提供物理真值，机器学习负责候选筛选，真实 GA 提供优化基准，最后用真实验证说明预测筛选有没有用。", 0.9, 1.48, 11.55, 0.62, 21, True)
    add_process(slide, ["物理真值计算\nCOMSOL 频散", "条件预测模型\n分类 + 回归", "真实 GA 优化\nCOMSOL overlap", "预测筛选对比\nTop5 / 随机 / GA20"], 0.9, 2.75, 11.55, 1.85)
    fit_image(slide, assets["ch1_route"], 1.05, 5.05, 10.95, 1.25, border=False)
    preview.append(("本文研究思路", "物理真值计算 → 条件预测建模 → 真实闭环优化 → 预测筛选对比。"))

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 5, "二维参数化声子晶体及目标频带任务")
    fit_image(slide, assets["unit"], 0.9, 1.55, 4.95, 4.65)
    add_textbox(slide, "任务输入/输出", 6.35, 1.6, 2.0, 0.4, 21, True)
    add_body_lines(slide, ["输入：结构参数、结构族信息、目标频带", "输出：目标频带是否有效开启、目标频带覆盖率", "同一结构在不同目标频带下可对应不同标签"], 6.35, 2.18, 5.75, 1.35, 17)
    add_textbox(slide, "六个目标频带", 6.35, 4.1, 2.0, 0.3, 16, True, COLORS["green"])
    add_target_bands(slide, 6.35, 4.55, 5.25, 1.7)
    preview.append(("研究对象与目标频带", "二维参数化单胞；六个目标频带；输入结构与目标条件，输出开启和覆盖率。"))

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 6, "基于 Bloch 周期边界的频散计算方法")
    add_process(slide, ["周期单胞", "Bloch-Floquet\n边界条件", "波矢路径\n特征频率", "频散曲线", "目标重叠\n评价"], 0.82, 1.74, 11.75, 1.35)
    add_formula_overlap(slide, 0.95, 3.62)
    fit_image(slide, assets["overlap"], 6.92, 3.25, 5.15, 2.85)
    add_textbox(slide, "这一页只讲计算链条和评价指标，不展开过多理论推导。", 1.0, 5.12, 5.45, 0.45, 16, False, COLORS["muted"])
    preview.append(("COMSOL 频散计算方法", "Bloch 周期边界；波矢路径求特征频率；overlap/coverage 评价。"))

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 7, "单胞几何模型与参数化描述")
    fit_image(slide, assets["geom_workflow"], 0.82, 1.55, 6.25, 4.9)
    add_card_text(slide, "结构族 + 连续参数", "候选结构由形状族和几何参数共同控制。", 7.55, 1.65, 4.25, 1.05, COLORS["green"])
    add_card_text(slide, "有效性检查", "几何有效性、接触有效性和求解稳定性决定能否进入 COMSOL。", 7.55, 3.05, 4.25, 1.05, COLORS["blue"])
    add_card_text(slide, "优化惩罚", "无效结构在真实 GA 中给予惩罚，避免搜索被失败个体牵引。", 7.55, 4.45, 4.25, 1.05, COLORS["orange"])
    preview.append(("几何参数化与候选生成", "结构族和参数控制候选；有效结构进入 COMSOL；无效结构惩罚。"))

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 8, "面向目标频带的条件分类与回归模型")
    fit_image(slide, assets["ml_model"], 0.8, 1.5, 5.55, 4.7)
    add_textbox(slide, "条件输入是关键", 6.9, 1.68, 2.4, 0.36, 21, True, COLORS["green"])
    add_body_lines(slide, ["分类模型：判断指定目标频带内是否形成有效带隙", "回归模型：预测目标频带重叠宽度或覆盖率", "把目标频带上下限作为输入，使同一结构可被不同频带条件重新评价"], 6.9, 2.28, 4.95, 1.55, 16)
    fit_image(slide, assets["ml_metrics"], 6.9, 4.25, 4.95, 1.8)
    preview.append(("机器学习条件预测模型", "分类判断开启；回归预测 overlap/coverage；目标频带上下限作为输入。"))

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 9, "预测模型的作用：候选筛选而非最终判据")
    add_textbox(slide, "预测模型用于快速初筛和排序；最终性能仍以 COMSOL 真实频散结果为准。", 0.9, 1.55, 11.5, 0.62, 24, True)
    add_card_text(slide, "机器学习预测", "速度快，适合大规模候选初筛；可信度依赖后续真实验证。", 1.0, 2.65, 5.15, 1.6, COLORS["green"])
    add_card_text(slide, "COMSOL 频散计算", "速度慢，但物理可信度高，是标签来源、GA 适应度和最终判据。", 7.1, 2.65, 5.15, 1.6, COLORS["blue"])
    add_round_rect(slide, 1.28, 5.35, 10.75, 0.75, RGBColor(252, 248, 239), RGBColor(230, 207, 149))
    add_textbox(slide, "答辩口径：预测排序本身不是终点；价值在于相同 COMSOL 预算下更快找到值得验证的结构。", 1.58, 5.55, 10.15, 0.28, 16, True, COLORS["gold"], PP_ALIGN.CENTER)
    preview.append(("模型作用定位", "预测用于初筛与排序；COMSOL 是最终判据；价值是节省验证预算。"))

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 10, "基于真实频散计算的目标频带遗传优化")
    fit_image(slide, assets["ga_flow"], 0.8, 1.5, 5.65, 4.65)
    fit_image(slide, assets["ga_comsol"], 6.75, 1.5, 5.35, 3.05)
    add_round_rect(slide, 6.95, 4.95, 4.95, 0.8, RGBColor(242, 247, 245), RGBColor(205, 224, 216))
    add_textbox(slide, "关键点：这里的 GA 适应度来自真实 COMSOL overlap，不是预测模型。", 7.2, 5.15, 4.45, 0.28, 15, True, COLORS["green"], PP_ALIGN.CENTER)
    preview.append(("真实 COMSOL-in-loop GA", "每个个体进入 COMSOL；overlap 作为适应度；六个频带分别优化。"))

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 11, "不同目标频带下的真实优化结果")
    fit_image(slide, assets["ga_best"], 0.75, 1.45, 5.85, 3.15)
    fit_image(slide, assets["ga_conv"], 6.9, 1.45, 5.35, 3.15)
    fit_image(slide, assets["ga_cells"], 0.95, 4.88, 5.2, 1.55, border=False)
    add_body_lines(slide, ["180-220 Hz 可实现完整覆盖", "160-200 Hz、200-240 Hz 覆盖率较高", "220-260 Hz、240-280 Hz 高频目标覆盖较低", "高频困难说明当前结构族和参数化空间表达能力有限"], 7.0, 4.88, 4.8, 1.5, 14)
    preview.append(("GA 优化总体结果", "20代真实GA：180-220完整覆盖，中频较好，高频覆盖较低。"))

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 12, "独立候选集验证：预测 Top5 vs 随机候选")
    fit_image(slide, assets["holdout_pipe"], 0.8, 1.48, 6.0, 4.85)
    add_textbox(slide, "严格独立验证", 7.35, 1.7, 2.5, 0.35, 21, True, COLORS["green"])
    add_body_lines(slide, ["重新构建未见候选集", "排除 v12 训练集样本", "排除第 4 章 GA20 历史记录", "排除已有第 5 章候选", "同一候选池内比较预测 Top5 与随机候选"], 7.35, 2.28, 4.75, 2.2, 16)
    add_round_rect(slide, 7.35, 5.35, 4.55, 0.65, RGBColor(242, 247, 245), RGBColor(205, 224, 216))
    add_textbox(slide, "每个频带 Top5 + random5，共 60 次 COMSOL 验证。", 7.58, 5.53, 4.05, 0.22, 14, True, COLORS["green"], PP_ALIGN.CENTER)
    preview.append(("预测筛选实验设置", "严格独立候选池；排除训练/GA20/已有候选；Top5 vs random5。"))

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 13, "预测筛选在多数频带下优于随机候选")
    fit_image(slide, assets["strict_overlap"], 0.75, 1.45, 5.75, 2.35)
    fit_image(slide, assets["strict_ratio"], 6.9, 1.45, 5.2, 2.35)
    add_textbox(slide, "严格 holdout 结果", 0.95, 4.25, 2.5, 0.35, 20, True, COLORS["green"])
    add_body_lines(slide, ["预测 Top5 在 5/6 个频带上的最优 overlap 高于随机候选", "180-220 Hz 达到 GA20 基准的 0.999", "200-240 Hz 达到 GA20 基准的 0.986", "240-280 Hz 中预测与随机均较弱，属于方法边界"], 0.95, 4.78, 5.55, 1.45, 15)
    fit_image(slide, assets["strict_cover"], 7.0, 4.35, 5.1, 1.85)
    preview.append(("Top5 与随机候选对比", "严格 holdout：预测Top5在5/6频带优于随机；180-220与200-240接近GA20。"))

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 14, "典型目标频带结构频散验证")
    fit_image(slide, assets["typical_cells"], 0.75, 1.48, 5.7, 4.65)
    fit_image(slide, assets["typical_disp_strict"], 6.75, 1.48, 5.65, 4.65)
    add_textbox(slide, "选取 180-220 Hz、200-240 Hz 和 240-280 Hz：分别对应正例、过渡频带和困难频带。", 1.05, 6.25, 11.0, 0.3, 15, True, COLORS["muted"], PP_ALIGN.CENTER)
    preview.append(("典型案例分析", "结构图 + 频散图；180-220正例，200-240过渡，240-280边界。"))

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 15, "高频目标频带的困难与方法边界")
    fit_image(slide, assets["highfreq"], 0.8, 1.55, 6.1, 4.9)
    add_textbox(slide, "讲法要稳", 7.35, 1.72, 2.2, 0.35, 21, True, COLORS["red"])
    add_body_lines(slide, ["220-260 Hz、240-280 Hz 的预测候选、随机候选和真实 GA 均表现有限", "问题不只是模型预测不准，也可能是当前结构族难以表达高频带隙", "高频补样能找到局部弱频带候选，但不等于完整目标带隙", "最终仍需 COMSOL 频散验证"], 7.35, 2.28, 4.75, 2.1, 16)
    add_round_rect(slide, 7.35, 5.55, 4.55, 0.58, RGBColor(252, 244, 242), RGBColor(228, 184, 177))
    add_textbox(slide, "高频弱，揭示的是当前结构族适用边界。", 7.58, 5.72, 4.05, 0.22, 14, True, COLORS["red"], PP_ALIGN.CENTER)
    preview.append(("高频困难频带分析", "高频表现弱；边界来自结构族与参数空间；最终仍需COMSOL验证。"))

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 16, "主要结论")
    conclusions = [
        ("COMSOL 真值流程", "建立二维周期声子晶体目标频带频散计算与评价指标体系。", COLORS["green"]),
        ("条件预测模型", "构建分类与回归模型，用于候选结构初筛和排序。", COLORS["blue"]),
        ("真实 GA 基准", "建立以 COMSOL overlap 为适应度的闭环 GA 优化基准。", COLORS["orange"]),
        ("严格验证结果", "独立验证表明预测 Top5 多数频带优于随机候选，提高有限预算下的发现效率。", COLORS["gold"]),
    ]
    for i, (title, body, color) in enumerate(conclusions):
        add_card_text(slide, title, body, 0.95 + (i % 2) * 5.85, 1.75 + (i // 2) * 2.05, 5.1, 1.45, color)
    add_round_rect(slide, 1.25, 6.0, 10.75, 0.58, RGBColor(242, 247, 245), RGBColor(205, 224, 216))
    add_textbox(slide, "机器学习预测模型不能替代 COMSOL，但可以提高 COMSOL 验证资源的使用效率。", 1.55, 6.15, 10.15, 0.22, 16, True, COLORS["green"], PP_ALIGN.CENTER)
    preview.append(("主要结论", "COMSOL流程；条件预测；真实GA；严格独立验证支持筛选价值。"))

    slide = prs.slides.add_slide(blank); set_slide_bg(slide); add_title(slide, 17, "不足与展望")
    add_card_text(slide, "结构族边界", "当前结构族和参数化空间对高频目标频带表达能力有限。", 1.0, 1.9, 3.45, 2.0, COLORS["red"])
    add_card_text(slide, "数据稀疏", "高频目标区间有效样本仍然稀疏，模型外推能力受限。", 4.95, 1.9, 3.45, 2.0, COLORS["orange"])
    add_card_text(slide, "后续方向", "引入更丰富结构族、主动采样策略和物理约束更强的预测模型。", 8.9, 1.9, 3.45, 2.0, COLORS["green"])
    add_textbox(slide, "后续研究将进一步提高高频目标频带样本覆盖度，并探索更具表达能力的结构参数化方法。", 1.28, 5.05, 10.85, 0.55, 23, True, COLORS["ink"], PP_ALIGN.CENTER)
    preview.append(("不足与展望", "高频结构表达能力有限；高频数据稀疏；丰富结构族和主动采样。"))

    slide = prs.slides.add_slide(blank); set_slide_bg(slide)
    cover_image(slide, assets["cover"], 0, 0, WIDE_W, WIDE_H, True)
    add_textbox(slide, "请各位老师批评指正", 2.1, 2.7, 8.9, 0.75, 42, True, COLORS["ink"], PP_ALIGN.CENTER)
    add_textbox(slide, "谢谢！", 5.55, 3.65, 2.2, 0.5, 25, True, COLORS["green"], PP_ALIGN.CENTER)
    preview.append(("致谢", "请各位老师批评指正。"))

    prs.save(str(OUT_PPTX))
    build_preview(preview)
    print(OUT_PPTX)
    print(CONTACT_SHEET)


if __name__ == "__main__":
    main()
