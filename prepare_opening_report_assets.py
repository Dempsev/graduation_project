from pathlib import Path
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from reportlab.lib.colors import HexColor
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfbase.pdfmetrics import registerFont
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer


ROOT = Path(r"D:\graduation_project\coad")
SRC_PPT = Path(r"D:\graduation_project\opening_report_continued_2026-03-25.pptx")
OUT_PPT = ROOT / "opening_report_continued_2026-03-25_cn.pptx"
OUT_PDF = ROOT / "opening_report_talk_notes_2026-04-02.pdf"


REPLACEMENTS = [
    ("seed-only refined / whitelist GA", "仅种子细化 / 白名单遗传算法"),
    ("primary + probe + diversity", "主选 + 探测 + 多样性"),
    ("post-ranking refinement", "排序后的局部精修"),
    ("seed vs GA local tuner", "种子方案与遗传算法局部微调对比"),
    ("College of Aerospace Engineering", "航空航天学院"),
    ("Stage1-Stage4", "第一到第四阶段"),
    ("Stage1-2", "第一、二阶段"),
    ("calibrated scoring", "校准后的评分"),
    ("calibration summary", "校准摘要"),
    ("contact gate", "接触门槛"),
    ("candidate pool", "候选池"),
    ("positive family", "正增益族群"),
    ("shape family", "形状族"),
    ("broad rollout", "大范围铺开探索"),
    ("diversity-aware", "兼顾多样性"),
    ("mean gain", "平均增益"),
    ("train-ready", "可训练"),
    ("dataset v1", "数据集 v1"),
    ("Stage2 main effects", "第二阶段主效应"),
    ("Stage2 refine", "第二阶段细化"),
    ("v11 calibration", "v11 校准"),
    ("rows_total", "总条数"),
    ("best gate", "最优门槛"),
    ("fixed 3-4", "固定 3-4"),
    ("CONTENTS", "目录"),
    ("Stage1", "第一阶段"),
    ("Stage2", "第二阶段"),
    ("Stage3", "第三阶段"),
    ("Stage4", "第四阶段"),
    ("validation", "验证"),
    ("baseline", "基线"),
    ("dataset", "数据集"),
    ("surrogate", "代理模型"),
    ("manifest", "候选清单"),
    ("shortlist", "短名单"),
    ("scoring", "评分"),
    ("whitelist", "白名单"),
    ("refined", "细化后"),
    ("refine", "细化"),
    ("refinement", "精修"),
    ("seed-only", "仅种子"),
    ("seed family", "种子族群"),
    ("seed", "种子"),
    ("family", "族群"),
    ("shape", "形状"),
    ("exploitation", "定向深挖"),
    ("rollout", "铺开探索"),
    ("broad", "大范围"),
    ("calibration", "校准"),
    ("database", "数据库"),
    ("candidate", "候选"),
    ("positive", "正增益"),
    ("precision / recall", "精确率 / 召回率"),
    ("precision", "精确率"),
    ("contact", "接触"),
    ("pool", "池"),
    ("point", "点位"),
    ("center", "中心"),
    ("contour", "轮廓"),
    ("rows", "条数"),
    ("count", "数量"),
    ("total", "总数"),
    ("gated", "筛后"),
    ("Harmonics", "高阶谐波项"),
    ("harmonics", "高阶谐波项"),
    ("Hz", "赫兹"),
    ("GA rows", "遗传算法样本数"),
    ("GA v1", "遗传算法 v1"),
    ("GA", "遗传算法"),
    ("FN", "假负例"),
    ("FP", "假正例"),
    ("TP", "真正例"),
]


SLIDE_NOTES = [
    (
        "封面",
        [
            "各位老师好，我汇报的题目是“贪吃蛇算法启发的声子晶体混合拓扑生成及数据驱动的逆向设计”。今天这次汇报主要是在开题基础上，补充我最近一段时间已经完成的新增进展。",
            "我这次想重点说明两件事：第一，原来只是设想的数据库部分，现在已经真正运转起来；第二，验证、评分和局部优化开始形成闭环，不再只是分散的单点实验。",
        ],
    ),
    (
        "目录页",
        [
            "这一页是整体结构。我先简要带一下研究背景和现状，然后重点汇报本阶段已经推进出来的工作，最后说明开题之后新增的实质性进展。",
            "如果老师更关心新增工作，可以重点看后半部分，尤其是数据库、真实验证和校准这三块。",
        ],
    ),
    (
        "汇报提纲",
        [
            "这页提纲和上一页一致，主要是提醒后面新增内容会占更大比重。因为和开题相比，当前最大的变化不是研究目标变了，而是执行路径更清楚了。",
            "也就是说，现在不是停留在“准备做什么”，而是已经进入“哪些环节跑通了、哪些问题暴露出来了、后面怎么继续收紧策略”。",
        ],
    ),
    (
        "数据库构建",
        [
            "这一页想表达的是，数据库构建现在已经不是附属工作，而是整个项目的中枢。当前我正在做持续验证，数据库也在从早期筛样逐步扩展到真值回灌。",
            "开题时数据库更多还是一个规划概念，但现在已经积累出结构化样本、可训练标签以及多轮验证记录。后面的代理模型、候选清单和校准，都会统一建立在这套口径一致的数据基础上。",
            "所以“继续补验证”本身并不只是补图补结果，它同时也是数据库构建的一部分。",
        ],
    ),
    (
        "四个阶段总览",
        [
            "这里我把当前工作分成四个阶段。第一阶段是受控筛样，核心是找可信母体和候选族群；第二阶段是参数规律提炼；第三阶段是把样本和标签整理成可训练数据库；第四阶段则是回到 COMSOL 做真实验证。",
            "这四个阶段不是线性做完一次就结束，而是第四阶段的真值结果会不断回灌到第三阶段数据库，再反过来影响评分和下一轮短名单。",
            "因此我现在更愿意把整个工作理解成一个循环系统，而不是一条单向流程。",
        ],
    ),
    (
        "第一阶段",
        [
            "第一阶段的任务不是直接找最终冠军结构，而是先判断哪些形状族真正值得继续深挖。因为如果母体不可信，后面的优化很容易建立在不稳定样本上。",
            "这一阶段的意义在于固定比较基线，并且从早期随机和规则样本里筛出正增益族群，作为后续深入搜索的种子。",
            "也正因为第一阶段之后，我已经识别出相对可信的种子族群，所以后面的策略开始从大范围撒网，收缩到围绕已知有效族群做定向深挖。",
        ],
    ),
    (
        "第二阶段概览",
        [
            "这一页是第二阶段的入口。核心思想是从“发现哪些族群值得做”，过渡到“这些族群里到底哪些参数真正在起作用”。",
            "也就是说，第二阶段开始不再追求盲目扩大搜索规模，而是拆解参数贡献，为后面的精修和排序提供依据。",
        ],
    ),
    (
        "第二阶段细化",
        [
            "在第二阶段里，我重点区分主效应、可细化区间和高阶项优先级。结果上看，像 a1、a2、b2 这类低阶参数，已经形成比较稳定的有效工作区间。",
            "而 a4、b5 这类高阶项更适合作为局部强化项，而不是一上来就大范围扫。这个判断直接影响后面搜索策略，因为它说明哪些参数值得先锁定，哪些参数只在局部精修时再放开。",
        ],
    ),
    (
        "第三阶段数据库",
        [
            "第三阶段本身不直接产生新的物理结果，但它决定了前面筛样结果和后面真实验证能不能被统一利用。我的目标是把样本表、标签口径、候选池和候选清单全部拉到同一套数据框架里。",
            "如果没有这一步，前面的实验只是分散在不同表格里的孤立结果，既无法稳定训练代理模型，也很难追踪短名单规则到底为什么有效或者失效。",
            "所以第三阶段的价值，本质上是把经验试验转成可持续积累、可反复调用的数据资产。",
        ],
    ),
    (
        "第四阶段真实验证",
        [
            "第四阶段的作用，是把短名单真正拉回 COMSOL 做真值验证。只有这一步回来，才能知道当前评分有没有偏、候选清单是不是过窄，以及哪些族群在真实条件下更稳定。",
            "这一阶段最关键的不是又多做了几次仿真，而是这些真值结果开始反向约束前面的评分和筛选策略，让整个流程真正闭环。",
        ],
    ),
    (
        "策略收敛过程",
        [
            "这一页可以理解成开题之后策略是怎么一步步收紧的。最早我还是偏大范围探索，希望先扩展几何空间；但很快发现接触条件和真实增益稳定性都不够。",
            "随后策略收缩到仅围绕种子族群做细化，再进一步引入兼顾多样性的候选清单，最后发展到校准加白名单后的局部遗传算法微调。",
            "这条变化线说明我不是简单追求更复杂的算法，而是让每一步都更贴近真实验证反馈。",
        ],
    ),
    (
        "遗传算法定位",
        [
            "这一页我想特别强调，遗传算法在当前阶段不是主搜索器，而是白名单之后的局部微调器。只对白名单形状开放，是为了避免搜索再次发散。",
            "从参数变化幅度也能看出来，它不是在创造全新规律，而是在可信点附近做小范围精修。因此它最合适的定位是排序之后的补充优化，而不是取代前面的筛样和校准。",
        ],
    ),
    (
        "真实验证结果",
        [
            "这页展示的是 v10 这一轮仅种子细化方案的真实验证结果。重点不只是分数高，而是短名单真正回到了物理仿真，而且有 8 个候选里 6 个取得正增益。",
            "这说明当前这条围绕种子族群细化的主线，已经具备稳定回收有效样本的能力。剩下暴露出来的主要问题，更多集中在接触门槛，而不是排序完全失效。",
        ],
    ),
    (
        "校准后的评分",
        [
            "有了前面的真值回证之后，v11 开始尝试校准后的评分。目的不是一味把样本砍得更狠，而是让短名单的精确率和召回率处在更稳定的折中位置。",
            "从当前结果看，策略明显更强调接触稳定性，先保证入选样本不要在接触条件上整体塌掉。这也说明数据库、真实验证和评分规则之间，已经开始互相约束而不是各做各的。",
        ],
    ),
    (
        "下一步计划",
        [
            "下一步我仍然沿着数据库、验证、校准这条主线推进。短期内会继续补一小轮真实验证，把新增样本持续回灌到数据库里。",
            "同时，我会围绕接触不稳定的族群做误差归因，区分到底是种子选择、点位选择，还是几何接触本身过于脆弱。主标签仍保持固定 3 到 4，不轻易改口径，保证数据库、候选清单和评分规则能持续迭代。",
        ],
    ),
    (
        "阶段性结论",
        [
            "如果只看新增部分，目前已经能比较明确地得出几个结论。第一，项目已经从“准备构建代理模型”推进到“有真实验证回灌的数据库闭环”。",
            "第二，第一、二阶段已经把值得继续做的族群和参数规律明显收缩出来。第三，当前最核心的瓶颈不再是有没有候选，而是接触门槛稳定性和短名单校准是否足够保守。",
            "第四，遗传算法会继续保留，但只作为白名单后的局部精修工具，不再承担新的主搜索任务。",
        ],
    ),
    (
        "结束页",
        [
            "我的汇报就到这里。整体上，我想向老师们说明的是，这个项目现在已经从前期探索，进入了以数据库闭环和真实验证为核心的稳步推进阶段。",
            "后面也欢迎老师们重点从模型可信性、验证设计和下一步执行节奏这几个方面给我建议。谢谢各位老师。",
        ],
    ),
]


POSSIBLE_QA = [
    (
        "1. 你这项工作的创新点到底是什么？",
        "可以回答为两层：第一层是把拓扑生成、数据库构建、真实验证和校准评分连成闭环，而不是只做单次筛样；第二层是把遗传算法降级为白名单后的局部精修器，让搜索流程更稳、更可解释。",
    ),
    (
        "2. 为什么数据库构建会被你放到这么核心的位置？",
        "因为现在真正影响后续代理模型、候选清单和阈值校准的，不是哪一次单独实验，而是所有样本能不能在同一口径下统一积累。如果没有数据库，后面的模型训练和规则更新都不可持续。",
    ),
    (
        "3. 为什么不直接用遗传算法做全局搜索？",
        "可以强调两点：一是全局遗传算法容易重新发散，二是当前问题的瓶颈更多在接触稳定性而不是搜索器不够强。因此把遗传算法限定在白名单后做局部微调，更符合目前阶段的需求。",
    ),
    (
        "4. 你怎么证明代理模型或评分规则是可信的？",
        "核心回答是“靠真实回证，不靠自我循环”。也就是短名单必须回到 COMSOL 做验证，然后把真值结果重新回灌数据库，再用这些真值去校准评分阈值。",
    ),
    (
        "5. 现在最主要的技术瓶颈是什么？",
        "建议直接答“接触门槛稳定性”。因为目前不是没有候选，而是部分候选在真实仿真中容易因为接触条件而失稳，这会影响短名单的保守性和验证回收率。",
    ),
    (
        "6. 你后面 2 到 3 周最具体要做什么？",
        "可以按三点回答：继续补一小轮真实验证；围绕失效样本做误差归因；在不改变主标签口径的前提下继续更新数据库和校准规则。",
    ),
    (
        "7. 为什么主标签保持固定 3 到 4？",
        "因为当前更需要保证口径稳定，否则数据库和评分规则会跟着来回变化。先保持标签定义稳定，才能判断新增验证到底带来了什么真实信息。",
    ),
    (
        "8. 如果老师问“项目风险在哪”，怎么答？",
        "可以说风险主要有两类：一类是接触不稳定导致高分样本不能稳定转化为真值增益；另一类是样本规模还需要继续扩充，避免校准结论过早定型。但这两类风险现在都已经进入可跟踪、可验证的状态。",
    ),
]


TERM_EXPLANATIONS = [
    (
        "1. 老师问“你这里说的种子是什么意思？”",
        "可以这样答：这里的“种子”不是算法里的随机种子，而是指在前期筛样里已经表现出正增益潜力、并且相对可信的一批初始结构或初始点。后面的细化搜索不是从零开始乱找，而是围绕这些已经被证明有希望的起点继续做局部探索。",
    ),
    (
        "2. 老师问“种子族群和单个种子有什么区别？”",
        "可以回答：单个种子更像一个具体样本点，种子族群则是一类相近结构的集合。我现在更关注族群，而不是某一个孤立最优点，因为单点可能偶然性较强，族群才更能说明某种结构规律是否稳定存在。",
    ),
    (
        "3. 老师问“白名单是什么意思？”",
        "可以答：白名单就是经过前面筛样、评分和经验判断后，被认定为值得继续优化的一小部分可信候选集合。后续像遗传算法这样的局部精修，只在这个集合里开放，避免重新回到全空间发散搜索。",
    ),
    (
        "4. 老师问“短名单和白名单有什么区别？”",
        "可以解释为：白名单更像允许继续深入的可信范围，规模相对大一些；短名单则是某一轮真正拿去做 COMSOL 真实验证的更小一批候选。也就是白名单偏“可进入下一步”，短名单偏“本轮马上验证”。",
    ),
    (
        "5. 老师问“接触门槛是什么意思？”",
        "可以这样说：接触门槛描述的是一个候选结构在几何接触或物理接触条件上是否足够稳定、足够可信。如果这个条件不过关，哪怕评分看起来高，真实仿真里也可能失效。所以我现在把它当成短名单稳定性的关键约束之一。",
    ),
    (
        "6. 老师问“校准是在校准什么？”",
        "可以回答：校准的对象主要是评分规则和筛选阈值。也就是利用已经拿到的真实验证结果，回头检查哪些分数段真的更容易得到正增益，哪些阈值会让精确率和召回率更平衡，而不是只相信模型内部排序。",
    ),
    (
        "7. 老师问“代理模型在你这里起什么作用？”",
        "可以答：代理模型不是最后裁决者，而是提高筛选效率的工具。它先帮助我在大量候选里做初步排序和压缩范围，但最终是否可信，还是要回到真实仿真验证，所以它服务于闭环，不替代闭环。",
    ),
    (
        "8. 老师问“为什么你一直强调闭环？”",
        "可以答：因为如果只有筛样，没有真实验证回证，那前面的高分结果可能只是自洽，不一定真实有效。我现在强调闭环，就是要求‘筛样、建库、评分、验证、再回灌’形成往复修正，这样结论才越来越可靠。",
    ),
    (
        "9. 老师问“你这里的参数规律具体指什么？”",
        "可以回答：主要是想搞清楚不同参数对结果的贡献强弱、有效区间和优先级。比如哪些低阶参数一动就明显影响结果，哪些高阶参数更适合留到后面做微调，这些都属于参数规律。",
    ),
    (
        "10. 老师问“为什么保留遗传算法但不让它当主线？”",
        "可以答：因为现在的问题已经不是完全找不到方向，而是需要在可信区域里更稳地做局部改进。遗传算法如果当主线，容易重新把搜索空间放得太大；但如果把它放在白名单之后，它就能发挥局部微调的优势。",
    ),
]


def replace_text(text: str) -> str:
    for old, new in sorted(REPLACEMENTS, key=lambda item: len(item[0]), reverse=True):
        text = text.replace(old, new)
    return text


def process_text_frame(text_frame) -> int:
    changed = 0
    for paragraph in text_frame.paragraphs:
        full_text = "".join(run.text for run in paragraph.runs)
        if not full_text:
            continue
        replaced = replace_text(full_text)
        if replaced != full_text and paragraph.runs:
            paragraph.runs[0].text = replaced
            for run in paragraph.runs[1:]:
                run.text = ""
            changed += 1
    return changed


def process_shape(shape) -> int:
    changed = 0
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            changed += process_shape(sub_shape)
    if getattr(shape, "has_text_frame", False):
        changed += process_text_frame(shape.text_frame)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                changed += process_text_frame(cell.text_frame)
    return changed


def localize_ppt() -> list[tuple[int, str]]:
    prs = Presentation(str(SRC_PPT))
    for slide in prs.slides:
        for shape in slide.shapes:
            process_shape(shape)
    prs.save(str(OUT_PPT))

    leftovers = []
    english_pattern = re.compile(r"[A-Za-z][A-Za-z0-9+\-./ ]*[A-Za-z0-9%)]|\b[A-Za-z]{2,}\b")
    reloaded = Presentation(str(OUT_PPT))
    for index, slide in enumerate(reloaded.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                for match in english_pattern.finditer(shape.text):
                    token = match.group(0).strip()
                    if len(token) >= 2:
                        leftovers.append((index, token))
    return sorted(set(leftovers))


def build_pdf() -> None:
    registerFont(UnicodeCIDFont("STSong-Light"))

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "TitleCN",
        parent=styles["Title"],
        fontName="STSong-Light",
        fontSize=20,
        leading=26,
        alignment=TA_CENTER,
        textColor=HexColor("#16324f"),
        spaceAfter=14,
    )
    heading_style = ParagraphStyle(
        "HeadingCN",
        parent=styles["Heading2"],
        fontName="STSong-Light",
        fontSize=14,
        leading=20,
        textColor=HexColor("#16324f"),
        spaceBefore=10,
        spaceAfter=8,
    )
    body_style = ParagraphStyle(
        "BodyCN",
        parent=styles["BodyText"],
        fontName="STSong-Light",
        fontSize=10.8,
        leading=18,
        textColor=HexColor("#222222"),
        wordWrap="CJK",
        spaceAfter=6,
    )
    tip_style = ParagraphStyle(
        "TipCN",
        parent=body_style,
        textColor=HexColor("#5b4636"),
        backColor=HexColor("#f5efe6"),
        borderPadding=6,
    )

    doc = SimpleDocTemplate(
        str(OUT_PDF),
        pagesize=A4,
        leftMargin=1.8 * cm,
        rightMargin=1.8 * cm,
        topMargin=1.8 * cm,
        bottomMargin=1.6 * cm,
        title="开题答辩讲稿与可能提问",
        author="Codex",
    )

    story = [
        Paragraph("开题答辩讲稿与可能提问", title_style),
        Paragraph("对应文件：opening_report_continued_2026-03-25.pptx", body_style),
        Paragraph(
            "使用建议：这份讲稿按页对应原始幻灯片顺序编写，语气尽量口语化。实际答辩时可以把每页控制在 30 秒到 90 秒之间，遇到老师追问时优先围绕“数据库闭环、真实验证、校准逻辑”三条主线回答。",
            tip_style,
        ),
        Spacer(1, 0.3 * cm),
    ]

    for idx, (title, paragraphs) in enumerate(SLIDE_NOTES, 1):
        story.append(Paragraph(f"第 {idx} 页：{title}", heading_style))
        for paragraph in paragraphs:
            story.append(Paragraph(paragraph, body_style))

    story.append(PageBreak())
    story.append(Paragraph("老师可能提问与答法", title_style))
    for question, answer in POSSIBLE_QA:
        story.append(Paragraph(question, heading_style))
        story.append(Paragraph(answer, body_style))

    story.append(PageBreak())
    story.append(Paragraph("术语细节追问答法", title_style))
    story.append(
        Paragraph(
            "这一部分专门准备给老师追问术语时使用。建议回答时先用一句话给定义，再补一句“它在我这个项目里具体起什么作用”，这样最稳。",
            tip_style,
        )
    )
    for question, answer in TERM_EXPLANATIONS:
        story.append(Paragraph(question, heading_style))
        story.append(Paragraph(answer, body_style))

    doc.build(story)


if __name__ == "__main__":
    leftovers = localize_ppt()
    build_pdf()
    print(f"PPT -> {OUT_PPT}")
    print(f"PDF -> {OUT_PDF}")
    print("Remaining English tokens:")
    for slide_no, token in leftovers:
        print(f"slide {slide_no}: {token}")
