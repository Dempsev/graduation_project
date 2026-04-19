from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SRC_PPT = Path(r"D:\graduation_project\opening_report_continued_2026-03-25.pptx")
OUT_DIR = Path(r"d:\graduation_project\coad\output\ppt")
OUT_PPT = OUT_DIR / "opening_report_continued_2026-03-25_updated.pptx"

NAVY = RGBColor(31, 78, 121)
BLUE = RGBColor(91, 155, 213)
LIGHT_BLUE = RGBColor(221, 235, 247)
GREEN = RGBColor(112, 173, 71)
LIGHT_GREEN = RGBColor(226, 239, 218)
ORANGE = RGBColor(237, 125, 49)
LIGHT_ORANGE = RGBColor(252, 228, 214)
GRAY = RGBColor(89, 89, 89)
LIGHT_GRAY = RGBColor(242, 242, 242)
WHITE = RGBColor(255, 255, 255)


def set_text(shape, text: str, font_size: float, *, bold: bool = False, color: RGBColor = GRAY,
             align=PP_ALIGN.LEFT, font_name: str = "Microsoft YaHei") -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for idx, line in enumerate(text.split("\n")):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_round_box(slide, x, y, w, h, text: str, *, fill: RGBColor, line: RGBColor,
                  font_size: float = 12, text_color: RGBColor = GRAY, bold: bool = False):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    set_text(shape, text, font_size, bold=bold, color=text_color, align=PP_ALIGN.CENTER)
    return shape


def add_chevron(slide, x, y, w, h, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def add_panel(slide, x, y, w, h, title: str, title_fill: RGBColor):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = WHITE
    panel.line.color.rgb = title_fill
    panel.line.width = Pt(1.5)

    title_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + Inches(0.12), y + Inches(0.08), w - Inches(0.24), Inches(0.44))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = title_fill
    title_box.line.color.rgb = title_fill
    set_text(title_box, title, 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return panel


def clone_header_from_slide(src_slide, dst_slide) -> None:
    for idx in [0, 1, 2]:
        el = deepcopy(src_slide.shapes[idx]._element)
        dst_slide.shapes._spTree.insert_element_before(el, "p:extLst")


def move_slide(prs: Presentation, old_index: int, new_index: int) -> None:
    slide_id_list = prs.slides._sldIdLst  # type: ignore[attr-defined]
    slide_el = slide_id_list[old_index]
    slide_id_list.remove(slide_el)
    slide_id_list.insert(new_index, slide_el)


def update_slide_5(prs: Presentation) -> None:
    slide = prs.slides[4]
    set_text(slide.shapes[4], "□ 当前项目结构：三条主线如何衔接", 24, bold=True, color=GRAY)
    set_text(slide.shapes[6], "Physical Data\nGeneration & Validation", 19, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    set_text(
        slide.shapes[7],
        "目标：形成统一物理数据资产\n输入：Stage1 / Stage2 / Harmonics /\nStage4 validation\n输出：可信样本、固定 3-4 标签、\n持续回灌的数据闭环",
        15,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    set_text(slide.shapes[9], "Prediction\nBranch", 19, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    set_text(
        slide.shapes[10],
        "目标：shape + parameters\n→ 固定 3-4 bandgap targets\n输出：单阶段回归 baseline\n与双阶段分类+回归主线",
        15,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    set_text(slide.shapes[12], "Optimization\nBranch", 19, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    set_text(
        slide.shapes[13],
        "目标：先低成本筛种子，\n再进入真实 COMSOL 闭环\n输出：candidate pool、局部优化、\nreal GA 与真实增益",
        15,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    set_text(slide.shapes[15], "Legacy / Baseline\nComparison", 19, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    set_text(
        slide.shapes[16],
        "定位：旧版 scoring+cascade\n和 global surrogate-GA\n保留作对照，不再作为当前主线",
        15,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    set_text(
        slide.shapes[20],
        "这四块不是线性“一次做完”的关系，而是 Stage4 的真实结果会持续回灌统一数据集，"
        "再反过来影响 prediction 的评估和 optimization 的候选筛选。",
        15,
        color=GRAY,
    )


def update_slide_15(prs: Presentation) -> None:
    slide = prs.slides[14]
    set_text(slide.shapes[4], "□ 下一步会沿着“统一数据集 - 筛选 - 真值验证”主线推进", 22, bold=True, color=GRAY)
    set_text(
        slide.shapes[5],
        "对应当前项目结构，可以把“数据库构建、前向预测、物理闭环优化”重新串成更清晰的执行顺序。",
        15,
        color=GRAY,
    )
    set_text(
        slide.shapes[9],
        "当前主线闭环示意",
        17,
        bold=True,
        color=GRAY,
    )
    set_text(slide.shapes[12], "统一物理数据集", 15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    set_text(slide.shapes[15], "前向预测评估\n(单阶段 baseline / 双阶段主线)", 14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    set_text(slide.shapes[18], "候选筛选与保守局部优化", 15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    set_text(slide.shapes[21], "真实 COMSOL 验证\n与结果回灌", 15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    set_text(
        slide.shapes[10],
        "继续补一小轮 seed-only refined / whitelist GA 的真实 validation，把新样本持续回灌进统一数据集。\n"
        "围绕 contact 不稳定 family 做误差归因，判断是 seed 选择、point 选择还是几何接触脆弱性。\n"
        "保持 fixed 3-4 主标签不变，让 prediction 评估、manifest 和 scoring 规则一起继续收敛。",
        14,
        color=GRAY,
    )
    set_text(
        slide.shapes[22],
        "和开题时相比，变化不在于研究目标变了，而在于项目已经形成了“物理数据 - 前向预测 - 物理闭环优化”"
        "相互支撑的主线结构。",
        13,
        color=GRAY,
    )


def add_structure_slide(prs: Presentation) -> None:
    blank = prs.slide_layouts[1]
    slide = prs.slides.add_slide(blank)
    clone_header_from_slide(prs.slides[4], slide)

    title = slide.shapes.add_textbox(Inches(0.42), Inches(0.72), Inches(12.0), Inches(0.55))
    set_text(title, "□ 项目总体结构：物理数据、前向预测与物理闭环优化", 23, bold=True, color=GRAY)

    subtitle = slide.shapes.add_textbox(Inches(0.42), Inches(1.18), Inches(12.0), Inches(0.35))
    set_text(
        subtitle,
        "当前仓库的主线可以按三部分理解：Physical Data Generation & Validation、Prediction Branch、Optimization Branch。",
        13,
        color=GRAY,
    )

    panel_y = Inches(1.6)
    panel_h = Inches(4.95)
    panel_w = Inches(4.08)
    gap = Inches(0.16)
    left_x = Inches(0.34)
    mid_x = left_x + panel_w + gap
    right_x = mid_x + panel_w + gap

    add_panel(slide, left_x, panel_y, panel_w, panel_h, "Physical Data Generation & Validation", NAVY)
    add_panel(slide, mid_x, panel_y, panel_w, panel_h, "Prediction Branch (Forward Modeling)", GREEN)
    add_panel(slide, right_x, panel_y, panel_w, panel_h, "Optimization Branch (Physics-in-the-loop)", ORANGE)

    box_w = Inches(3.45)
    box_h = Inches(0.52)
    arrow_w = Inches(0.36)
    arrow_h = Inches(0.18)

    def col_x(base_x):
        return base_x + Inches(0.32)

    y0 = Inches(2.15)
    phys_items = [
        ("Stage1\nRandom snake-shape screening", LIGHT_BLUE, NAVY),
        ("Stage2\nLow-order sweep & local refinement", LIGHT_BLUE, NAVY),
        ("Stage2 (Harmonics)\nHigh-order exploration & refinement", LIGHT_BLUE, NAVY),
        ("Stage4\nPhysical validation & feedback", LIGHT_BLUE, NAVY),
        ("Unified Physical Dataset", NAVY, WHITE),
    ]
    for idx, (text, fill, line) in enumerate(phys_items):
        y = y0 + idx * Inches(0.78)
        add_round_box(slide, col_x(left_x), y, box_w, box_h, text, fill=fill, line=line, font_size=11.5, text_color=WHITE if idx == 4 else GRAY, bold=idx == 4)
        if idx < len(phys_items) - 1:
            add_chevron(slide, col_x(left_x) + Inches(1.55), y + Inches(0.56), arrow_w, arrow_h, BLUE)

    pred_items = [
        ("Task\nshape descriptors + parameters\n→ fixed 3-4 targets", LIGHT_GREEN, GREEN),
        ("Data filtering & targets\nvalid geometry/contact/solve\n gap34_Hz / width / is_open", LIGHT_GREEN, GREEN),
        ("Models\nsingle-stage regression baseline\n two-stage classification + regression", LIGHT_GREEN, GREEN),
        ("Evaluation\nMAE / RMSE / R²\nincluding Stage4 holdout", GREEN, WHITE),
    ]
    for idx, (text, fill, line) in enumerate(pred_items):
        y = y0 + idx * Inches(0.98)
        add_round_box(slide, col_x(mid_x), y, box_w, Inches(0.68), text, fill=fill, line=line, font_size=11.2, text_color=WHITE if idx == 3 else GRAY, bold=idx in {2, 3})
        if idx < len(pred_items) - 1:
            add_chevron(slide, col_x(mid_x) + Inches(1.55), y + Inches(0.72), arrow_w, arrow_h, GREEN)

    opt_items = [
        ("Front-end candidate filtering\n2 classifiers + 1 regressor", LIGHT_ORANGE, ORANGE),
        ("Strategy A\nconservative local parameter search\non fixed seed shapes", LIGHT_ORANGE, ORANGE),
        ("Limited COMSOL validation\nseed-upside assessment", LIGHT_ORANGE, ORANGE),
        ("Final optimization\nreal COMSOL-in-the-loop GA", LIGHT_ORANGE, ORANGE),
        ("Output\nvalidated optimal shape\n+ parameters + gain", ORANGE, WHITE),
    ]
    for idx, (text, fill, line) in enumerate(opt_items):
        y = y0 + idx * Inches(0.78)
        add_round_box(slide, col_x(right_x), y, box_w, box_h if idx != 4 else Inches(0.6), text, fill=fill, line=line, font_size=11.1, text_color=WHITE if idx == 4 else GRAY, bold=idx >= 3)
        if idx < len(opt_items) - 1:
            add_chevron(slide, col_x(right_x) + Inches(1.55), y + Inches(0.56), arrow_w, arrow_h, ORANGE)

    bottom = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.52), Inches(6.72), Inches(12.18), Inches(0.46))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = LIGHT_GRAY
    bottom.line.color.rgb = RGBColor(191, 191, 191)
    set_text(
        bottom,
        "Baseline comparison: legacy scoring + cascade surrogate, and global surrogate-GA are retained as comparison paths rather than the current thesis mainline.",
        11.5,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )

    line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.08), Inches(5.7), Inches(5.05), Inches(5.7))
    line1.line.color.rgb = BLUE
    line1.line.width = Pt(1.6)
    line2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.38), Inches(5.85), Inches(8.74), Inches(5.85))
    line2.line.color.rgb = ORANGE
    line2.line.width = Pt(1.6)

    move_slide(prs, len(prs.slides) - 1, 5)


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(SRC_PPT))
    update_slide_5(prs)
    update_slide_15(prs)
    add_structure_slide(prs)
    prs.save(str(OUT_PPT))
    print(f"[DONE] saved updated deck to {OUT_PPT}")


if __name__ == "__main__":
    main()
